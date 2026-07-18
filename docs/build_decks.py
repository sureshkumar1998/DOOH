#!/usr/bin/env python3
"""Generate two presentation decks (navy/gold theme) for the SunMobility Ad Network.

Model 1 — DOOH rationale (why DOOH, pros/cons, why not CDN+agency).
Model 2 — Investor pitch.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1E, 0x3A, 0x5F)
GOLD = RGBColor(0xD9, 0xA5, 0x2B)
DARK = RGBColor(0x22, 0x2B, 0x38)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY2 = RGBColor(0x16, 0x2D, 0x4A)

SW, SH = Inches(13.333), Inches(7.5)


def new_prs():
    p = Presentation()
    p.slide_width = SW
    p.slide_height = SH
    return p


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        if isinstance(para, tuple):
            para = [para]
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = 'Calibri'
    return tb


def footer(slide, n, title):
    text(slide, Inches(0.5), Inches(7.0), Inches(8), Inches(0.4),
         [[(f'SunMobility Ad Network — {title}', 9, GRAY, False)]])
    text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
         [[(str(n), 9, GRAY, False)]], align=PP_ALIGN.RIGHT)


def header(slide, kicker, title):
    rect(slide, 0, 0, SW, Inches(1.5), NAVY)
    rect(slide, 0, Inches(1.5), SW, Pt(4), GOLD)
    if kicker:
        text(slide, Inches(0.6), Inches(0.28), Inches(12), Inches(0.4),
             [[(kicker.upper(), 12, GOLD, True)]])
    text(slide, Inches(0.6), Inches(0.62), Inches(12.1), Inches(0.8),
         [[(title, 28, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)


def bullets(slide, items, top=Inches(1.9), left=Inches(0.7), width=Inches(12), size=16, gap=10):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        prefix = '     ◦  ' if lvl else '•  '
        col = GRAY if lvl else DARK
        runs.append([(prefix, size, GOLD, True), (txt, size, col, False)])
    text(slide, left, top, width, Inches(5), runs, space=gap)


def note(slide, txt, top=Inches(6.45)):
    rect(slide, Inches(0.7), top, Inches(11.9), Pt(2), LIGHT)
    text(slide, Inches(0.7), top + Pt(6), Inches(11.9), Inches(0.5),
         [[('▶  ', 12, GOLD, True), (txt, 12, GRAY, False)]])


def table(slide, rows, top, left=Inches(0.7), width=Inches(11.9), col_w=None,
          header_fill=NAVY, height=None):
    nrows, ncols = len(rows), len(rows[0])
    h = height or Inches(0.5 * nrows)
    gtbl = slide.shapes.add_table(nrows, ncols, left, top, width, h).table
    if col_w:
        for i, w in enumerate(col_w):
            gtbl.columns[i].width = w
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = gtbl.cell(ri, ci)
            c.margin_left = Inches(0.12); c.margin_right = Inches(0.12)
            c.margin_top = Inches(0.05); c.margin_bottom = Inches(0.05)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
                fg, bold, sz = WHITE, True, 13
            else:
                c.fill.solid(); c.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                fg, bold, sz = DARK, (ci == 0), 12
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.size = Pt(sz); r.font.color.rgb = fg; r.font.bold = bold; r.font.name = 'Calibri'
    return gtbl


def title_slide(prs, kicker, title, subtitle, tag):
    s = blank(prs)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, Inches(0.7), Inches(2.4), Inches(1.6), Pt(5), GOLD)
    text(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.5), [[(kicker.upper(), 14, GOLD, True)]])
    text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(1.6), [[(title, 46, WHITE, True)]])
    text(s, Inches(0.7), Inches(4.3), Inches(11.5), Inches(1.0), [[(subtitle, 20, RGBColor(0xC7,0xD2,0xE0), False)]])
    text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5), [[(tag, 13, GOLD, True)]])
    return s


# =====================================================================
# MODEL 1 — DOOH rationale
# =====================================================================
def build_model1(path):
    prs = new_prs()

    title_slide(prs, 'Advertising Strategy', 'SunMobility Ad Network',
                'Turning every battery swap into an advertising impression',
                'Model 1 · Why DOOH on our own swap-station network')

    s = blank(prs); header(s, 'Slide 2', 'The Opportunity')
    bullets(s, [
        'We operate 1,790+ swap stations across 23 zones (Bengaluru, Delhi, Jaipur, Hyderabad…)',
        'Every station already has a screen the customer looks at during a swap',
        'Each swap = a guaranteed, face-to-screen moment — a captive audience we already own',
        'Today, that screen-time is completely unmonetized',
    ])
    note(s, "We're sitting on a daily, recurring audience that nobody is currently selling.")
    footer(s, 2, 'Model 1')

    s = blank(prs); header(s, 'Slide 3', 'What is DOOH?')
    bullets(s, [
        'DOOH = Digital Out-Of-Home advertising — digital screens in physical public spaces',
        'Standard DOOH: malls, airports, transit, retail — viewers are passers-by',
        'Our version: the screen is on our own swap stations',
        ('The viewer is a customer mid-transaction, not a passer-by', 1),
        ('We control the hardware, network, schedule, and the data', 1),
    ])
    note(s, 'Standard DOOH hopes people glance. Ours sits in front of a captive user.')
    footer(s, 3, 'Model 1')

    # 3-layer visual
    s = blank(prs); header(s, 'Slide 4', 'Our Unique Asset: the captive 3-layer screen')
    cx, cw, ch = Inches(0.9), Inches(6.2), Inches(1.05)
    rect(s, cx, Inches(2.0), cw, ch, GOLD)
    text(s, cx, Inches(2.0), cw, ch, [[('TOP  →  ADVERTISEMENT', 16, NAVY, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, cx, Inches(3.15), cw, ch, NAVY)
    text(s, cx, Inches(3.15), cw, ch, [[('MIDDLE  →  SWAP UI', 16, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, cx, Inches(4.3), cw, ch, GOLD)
    text(s, cx, Inches(4.3), cw, ch, [[('BOTTOM  →  ADVERTISEMENT', 16, NAVY, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(7.4), Inches(2.0), Inches(5.3), Inches(4),
         [[('Why it works', 18, NAVY, True)],
          [('•  ', 15, GOLD, True), ('The customer stares at the middle layer to complete the swap', 15, DARK, False)],
          [('•  ', 15, GOLD, True), ('Ads (top + bottom) are in their field of view by design', 15, DARK, False)],
          [('•  ', 15, GOLD, True), ('Dwell time: 1–4 minutes per swap (1 / 2 / 3 battery packs)', 15, DARK, False)],
          [('•  ', 15, GOLD, True), ('Visibility ≈ 100% — not "maybe they looked"', 15, DARK, False)]], space=12)
    footer(s, 4, 'Model 1')

    s = blank(prs); header(s, 'Slide 5', 'Why We Chose DOOH (strategic fit)')
    bullets(s, [
        'We already own the screens — zero new placement cost',
        'Captive, high-intent audience — EV riders, daily commuters, fleet drivers',
        'Hyper-local targeting — down to a single station or a whole city zone',
        'Recurring footfall — the same customers swap daily, repeated exposure',
        'Real measurement — we have swap data, so impressions are counted, not guessed',
    ])
    note(s, 'The hardest parts of DOOH — screens, audience, location — we already have.')
    footer(s, 5, 'Model 1')

    s = blank(prs); header(s, 'Slide 6', 'DOOH: Pros')
    table(s, [
        ['Advantage', 'Why it matters for us'],
        ['Own the inventory', 'No rent to screen owners; ~100% margin on our network'],
        ['Captive attention', 'Swapper faces the screen — visibility ≈ 100%'],
        ['Local + zonal targeting', 'Sell a single station, a city, or the whole fleet'],
        ['Brand-safe & premium', 'No fraud, no bots, no ad-blockers — real screen, real person'],
        ['Data we own', 'Swap counts → defensible impression numbers'],
        ['Repeat exposure', 'Daily swappers see campaigns multiple times'],
    ], top=Inches(1.8), col_w=[Inches(3.6), Inches(8.3)])
    footer(s, 6, 'Model 1')

    s = blank(prs); header(s, 'Slide 7', 'DOOH: Cons (and how we mitigate)')
    table(s, [
        ['Limitation', 'Our mitigation'],
        ['Limited screen space (top/bottom only)', 'Cap ads at ≤60s; rotate by slot & schedule'],
        ['Audience capped by footfall', '1,790 stations + growing — scale by network size'],
        ['No clicks / direct response', 'Sell on reach & branding; add QR codes for response'],
        ['Content operations overhead', 'We built the dashboard: upload, target, schedule, preview'],
        ['Proving value to advertisers', 'Swap-data impressions — auditable, not estimated'],
    ], top=Inches(1.8), col_w=[Inches(4.6), Inches(7.3)])
    note(s, 'Each weakness is handled by the platform we already built.')
    footer(s, 7, 'Model 1')

    s = blank(prs); header(s, 'Slide 8', 'The Alternative: CDN + Ad Agency')
    bullets(s, [
        'Model: outsource — push content via a third-party CDN, let an ad agency / network sell & serve',
        'Agency brings advertisers and ad-serving technology',
        'CDN distributes creative to the endpoints',
        'We take a revenue share',
    ])
    note(s, 'The "easy" route: let someone else run the ad business and pay us a cut.')
    footer(s, 8, 'Model 1')

    s = blank(prs); header(s, 'Slide 9', 'CDN + Ad Agency: Pros')
    bullets(s, [
        'Fast to start — no platform to build; agency has advertiser relationships',
        'Outsourced sales — agency finds the buyers',
        'Mature ad-serving tech — targeting, rotation, reporting out of the box',
        'Programmatic demand — access to large ad exchanges / fill rates',
    ])
    note(s, 'Lower upfront effort — you plug in and collect a share.')
    footer(s, 9, 'Model 1')

    s = blank(prs); header(s, 'Slide 10', 'CDN + Ad Agency: Cons (why it does not fit us)')
    table(s, [
        ['Problem', 'Impact'],
        ['Revenue share', 'Agency + CDN take 30–50%+ — we keep a fraction of our own audience'],
        ['We lose the data', 'Agency owns impressions/reporting; we cannot prove our captive value'],
        ['Generic measurement', 'Online CPM / viewability ignores our swap-dwell advantage'],
        ['Brand control risk', 'Third-party ads could clash with our EV / brand context'],
        ['Wrong model', 'CDN/agency is built for web & app ads, not a captive swap screen'],
        ['Dependency', 'Monetization hostage to an external partner’s priorities'],
    ], top=Inches(1.8), col_w=[Inches(3.4), Inches(8.5)])
    footer(s, 10, 'Model 1')

    s = blank(prs); header(s, 'Slide 11', 'Side-by-Side')
    table(s, [
        ['', 'Own DOOH (chosen)', 'CDN + Ad Agency'],
        ['Setup effort', 'Higher — already built ✔', 'Low'],
        ['Revenue kept', '~100%', '50–70%'],
        ['Audience data', 'We own it', 'Agency owns it'],
        ['Measurement', 'Swap-driven, real', 'Generic CPM'],
        ['Brand control', 'Full', 'Limited'],
        ['Best for', 'Captive physical screens', 'Web / app inventory'],
    ], top=Inches(1.8), col_w=[Inches(3.0), Inches(4.7), Inches(4.2)])
    note(s, 'For a network we own with a captive audience, in-house wins on every axis that matters.')
    footer(s, 11, 'Model 1')

    s = blank(prs); header(s, 'Slide 12', 'Why In-House DOOH Wins for Us')
    bullets(s, [
        'We own the audience — keep the full value, don’t rent it out',
        'We own the data — swaps → impressions nobody else can measure',
        'Our screen is captive, not incidental — a premium we can charge for',
        'We already built the platform: upload, target (zone/station), schedule, preview, monitor, proof-of-play logs',
    ])
    note(s, 'The build cost is behind us; the margin advantage compounds forever.')
    footer(s, 12, 'Model 1')

    s = blank(prs); header(s, 'Slide 13', 'Our Differentiator: impressions from real swaps')
    bullets(s, [
        'No "visibility factor" guesswork like billboards',
        'impressions = total swap dwell-time ÷ ad loop length',
        'Driven by actual swap counts and battery-pack count per station, per day',
        'Auto-scales: busy day = more impressions, no manual tuning',
        'Advertiser sees: swaps · avg time-at-screen · guaranteed full plays · impressions',
    ])
    note(s, "We don't estimate eyeballs — we count swaps, and every swap is a face on the screen.")
    footer(s, 13, 'Model 1')

    s = blank(prs); header(s, 'Slide 14', 'Closing / Next Steps')
    bullets(s, [
        'A fully owned, measurable, captive-audience ad network across 1,790+ stations',
        'Platform live: upload → target by zone/station → schedule → preview → monitor → prove',
        'Next: onboard pilot advertisers in 2–3 high-footfall zones; report real impressions',
        'Ask: pilot approval / budget / advertiser introductions',
    ])
    note(s, 'We have the screens, the audience, the platform, and the proof — let’s start selling.')
    footer(s, 14, 'Model 1')

    prs.save(path)
    return len(prs.slides._sldIdLst)


# =====================================================================
# MODEL 2 — Investor pitch
# =====================================================================
def metric_row(slide, items, top=Inches(2.1)):
    n = len(items); gap = Inches(0.3)
    total_w = Inches(12.0); cw = Emu(int((total_w - gap * (n - 1)) / n))
    x = Inches(0.66)
    for (big, small) in items:
        rect(slide, x, top, cw, Inches(2.0), NAVY)
        rect(slide, x, top, cw, Pt(5), GOLD)
        text(slide, x, top + Inches(0.45), cw, Inches(0.9), [[(big, 40, GOLD, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, x, top + Inches(1.25), cw, Inches(0.6), [[(small, 13, WHITE, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        x = Emu(x + cw + gap)


def build_model2(path):
    prs = new_prs()

    title_slide(prs, 'Investor Pitch', 'SunMobility Ad Network',
                'A captive-audience advertising network built on India’s EV battery-swap infrastructure',
                'Model 2 · Investor Pitch')

    s = blank(prs); header(s, '01', 'The Problem & Opportunity')
    bullets(s, [
        'EV battery-swap stations are high-traffic, daily-use touchpoints — but their screen time earns nothing',
        'Brands struggle to reach urban EV riders, gig & fleet drivers with measurable, brand-safe media',
        'Online ads suffer fraud, ad-blockers and "did anyone really see it?" doubt',
        'Billboards/DOOH can’t prove who actually looked',
        'There’s an unmonetized, captive, measurable audience hiding in plain sight',
    ])
    note(s, 'A premium audience nobody is selling, on infrastructure we already operate.')
    footer(s, 2, 'Model 2')

    s = blank(prs); header(s, '02', 'The Asset: a captive network we already own')
    metric_row(s, [('1,790+', 'Swap stations live'), ('23', 'City zones'), ('1–4 min', 'Captive dwell / swap'), ('100%', 'Inventory owned')])
    text(s, Inches(0.66), Inches(4.5), Inches(12), Inches(1.5),
         [[('Every swap puts a rider in front of our screen for minutes — a recurring, ', 15, DARK, False),
           ('face-to-screen', 15, NAVY, True), (' audience that grows as the EV network grows.', 15, DARK, False)]])
    note(s, 'We don’t have to build an audience — we already serve one every day.')
    footer(s, 3, 'Model 2')

    s = blank(prs); header(s, '03', 'The Solution: an owned, measurable ad platform')
    bullets(s, [
        'A 3-layer screen: swap UI in the middle, ads top & bottom — captive by design',
        'A full ad-management dashboard: upload → target (zone/station) → schedule → preview',
        'Live fleet monitoring + proof-of-play logs (which ad ran, where, when)',
        'Impressions computed from real swap data — not estimated',
        'Built and running today on the existing station network',
    ])
    note(s, 'The product is live — this is execution, not a concept.')
    footer(s, 4, 'Model 2')

    s = blank(prs); header(s, '04', 'Why Now')
    bullets(s, [
        'EV two/three-wheeler adoption is scaling fast in India — swap volume rising',
        'DOOH is the fastest-growing ad medium; advertisers want brand-safe, measurable screens',
        'Retail-media & "own your audience" models are booming globally',
        'Our station network is already deployed — first-mover on swap-screen media',
    ])
    note(s, 'EV growth + DOOH growth + an owned network = a rare timing window.')
    footer(s, 5, 'Model 2')

    s = blank(prs); header(s, '05', 'How It Works')
    bullets(s, [
        'Advertiser picks a zone or specific stations and a campaign schedule',
        'Creative (≤60s) is uploaded, placed in a screen slot, previewed on a real CCU layout',
        'Stations pull content automatically; the player rotates ads around the swap UI',
        'Every swap’s dwell time → impressions; proof-of-play confirms delivery',
        'Advertiser gets a clear report: swaps · time-at-screen · full plays · impressions',
    ])
    footer(s, 6, 'Model 2')

    s = blank(prs); header(s, '06', 'Differentiator: impressions from real swaps')
    text(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(0.9),
         [[('impressions = total swap dwell-time  ÷  ad loop length', 22, NAVY, True)]], align=PP_ALIGN.CENTER)
    bullets(s, [
        'No billboard-style "visibility factor" guesswork',
        'Driven by actual swap counts and battery-pack count (1/2/3) per station, per day',
        'Captive screen → visibility ≈ 100%; auto-scales with footfall',
        'Auditable & defensible — we show the inputs, not a mystery CPM',
    ], top=Inches(3.0))
    note(s, 'A measurement advantage no billboard or web-ad network can match.')
    footer(s, 7, 'Model 2')

    s = blank(prs); header(s, '07', 'Business Model')
    table(s, [
        ['Lever', 'Detail'],
        ['Own the inventory', 'No screen rent — margin stays in-house (~100% of media revenue)'],
        ['Sell by geography', 'Single station · city zone · whole fleet — flexible packages'],
        ['Pricing', 'Premium CPM justified by captive, measurable, brand-safe attention'],
        ['Scales with network', 'Every new swap station adds ad inventory at ~zero extra cost'],
        ['Recurring demand', 'Campaign renewals from proven, reported impressions'],
    ], top=Inches(1.8), col_w=[Inches(3.4), Inches(8.5)])
    footer(s, 8, 'Model 2')

    s = blank(prs); header(s, '08', 'Why In-House vs. Outsourcing')
    table(s, [
        ['', 'Own platform (us)', 'CDN + Ad Agency'],
        ['Revenue kept', '~100%', '50–70%'],
        ['Audience data', 'We own it', 'Agency owns it'],
        ['Measurement', 'Swap-driven, real', 'Generic CPM'],
        ['Brand control', 'Full', 'Limited'],
        ['Fit for captive screen', 'Purpose-built', 'Built for web/app'],
    ], top=Inches(1.9), col_w=[Inches(3.4), Inches(4.4), Inches(4.1)])
    note(s, 'We keep the margin, the data, and the measurement edge by owning the stack.')
    footer(s, 9, 'Model 2')

    s = blank(prs); header(s, '09', 'Traction: the platform is built')
    bullets(s, [
        'Live integration with the station fleet API (1,790+ stations, 23 zones)',
        'Auto-authenticating, self-refreshing — no manual ops',
        'Ad dashboard: upload, 8:3 spec enforcement, zone/station targeting, scheduling, panel placement, preview',
        'Fleet monitor (online/offline + ads-running status) and proof-of-play logs',
        'Impressions engine validated against real swap data',
    ])
    note(s, 'De-risked: the hard engineering is done and running.')
    footer(s, 10, 'Model 2')

    s = blank(prs); header(s, '10', 'Illustrative Economics')
    bullets(s, [
        'A 130-swap/day station, one 60s ad: ≈ 210 impressions/day — every swapper sees ≥1 full play',
        'Scales by station count: hundreds of stations × daily swaps = millions of monthly impressions',
        'Inventory grows for free with each new swap station deployed',
        'Premium CPM × owned inventory × ~100% margin = high-leverage revenue',
    ])
    note(s, 'Illustrative — final numbers from pilot-zone measurement.')
    footer(s, 11, 'Model 2')

    s = blank(prs); header(s, '11', 'Go-To-Market')
    bullets(s, [
        'Phase 1 — Pilot: 2–3 high-footfall zones, onboard anchor advertisers, report real impressions',
        'Phase 2 — Productize: self-serve campaign booking by zone/station',
        'Phase 3 — Scale: roll out fleet-wide; pursue agency & brand direct deals on proven data',
        'Land-and-expand: pilot results → renewals & larger geographies',
    ])
    footer(s, 12, 'Model 2')

    s = blank(prs)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, Inches(0.7), Inches(2.3), Inches(1.6), Pt(5), GOLD)
    text(s, Inches(0.7), Inches(2.6), Inches(12), Inches(1.2), [[('The Ask & Vision', 40, WHITE, True)]])
    text(s, Inches(0.7), Inches(3.9), Inches(11.8), Inches(2.5),
         [[('•  ', 18, GOLD, True), ('Approve a pilot across 2–3 zones to prove revenue per station', 18, RGBColor(0xE5,0xEA,0xF1), False)],
          [('•  ', 18, GOLD, True), ('Vision: India’s largest captive, measurable EV-rider ad network', 18, RGBColor(0xE5,0xEA,0xF1), False)],
          [('•  ', 18, GOLD, True), ('Own the screens, own the data, own the margin', 18, RGBColor(0xE5,0xEA,0xF1), False)]], space=14)
    text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5), [[('SunMobility Ad Network', 13, GOLD, True)]])
    footer(s, 13, 'Model 2')

    prs.save(path)
    return len(prs.slides._sldIdLst)


if __name__ == '__main__':
    import os
    out = os.path.dirname(os.path.abspath(__file__))
    n1 = build_model1(os.path.join(out, 'SunMobility_Ad_Network_Model1_DOOH.pptx'))
    n2 = build_model2(os.path.join(out, 'SunMobility_Ad_Network_Model2_Investor.pptx'))
    print(f'Model 1 (DOOH rationale): {n1} slides')
    print(f'Model 2 (Investor pitch): {n2} slides')
